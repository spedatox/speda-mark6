import logging
import os
import re
import uuid
from functools import lru_cache
from pathlib import Path

from app.config import settings
from app.core.context import AgentContext
from app.skills.base import Skill

logger = logging.getLogger(__name__)

# ── PDF fonts ───────────────────────────────────────────────────────────────
# ReportLab's built-in Helvetica is CP1252-encoded and has no glyphs for the
# Turkish-specific letters ğ ş ı İ (they render as ▮). We bundle the DejaVu Sans
# family (full Unicode, permissive licence) and register it so PDFs render
# Turkish — and any Latin/Cyrillic/Greek — correctly on every deployment,
# independent of system fonts (Contabo container ships only what we bundle).
_FONTS_DIR = Path(__file__).parent / "fonts"


@lru_cache(maxsize=1)
def _register_pdf_fonts() -> dict[str, str]:
    """
    Register the bundled DejaVu family with ReportLab once per process and
    return the {role: font-name} map the PDF styles use. On any failure (fonts
    missing), fall back to the Helvetica family so generation never breaks —
    Turkish glyphs are lost but the document is still produced.
    """
    from reportlab.pdfbase import pdfmetrics            # type: ignore[import]
    from reportlab.pdfbase.ttfonts import TTFont        # type: ignore[import]

    fallback = {"normal": "Helvetica", "bold": "Helvetica-Bold",
                "italic": "Helvetica-Oblique", "mono": "Courier"}
    try:
        variants = {
            "DejaVuSans": "DejaVuSans.ttf",
            "DejaVuSans-Bold": "DejaVuSans-Bold.ttf",
            "DejaVuSans-Oblique": "DejaVuSans-Oblique.ttf",
            "DejaVuSans-BoldOblique": "DejaVuSans-BoldOblique.ttf",
            "DejaVuSansMono": "DejaVuSansMono.ttf",
        }
        for name, fname in variants.items():
            pdfmetrics.registerFont(TTFont(name, str(_FONTS_DIR / fname)))
        # Family mapping lets inline <b>/<i> markup resolve to the right variant.
        pdfmetrics.registerFontFamily(
            "DejaVuSans",
            normal="DejaVuSans", bold="DejaVuSans-Bold",
            italic="DejaVuSans-Oblique", boldItalic="DejaVuSans-BoldOblique",
        )
        return {"normal": "DejaVuSans", "bold": "DejaVuSans-Bold",
                "italic": "DejaVuSans-Oblique", "mono": "DejaVuSansMono"}
    except Exception as e:
        logger.warning(
            "pdf_font_register_failed",
            extra={"error": str(e), "fonts_dir": str(_FONTS_DIR)},
        )
        return fallback


# ── Markdown parsing ──────────────────────────────────────────────────────────

def _is_hr(line: str) -> bool:
    """A thematic break: 3+ of -, *, or _ (no pipes, so table separators are excluded)."""
    s = line.strip().replace(" ", "")
    return len(s) >= 3 and (set(s) == {"-"} or set(s) == {"*"} or set(s) == {"_"})


def _is_table_sep(line: str) -> bool:
    """A table separator row, e.g. |---|:--:|---| or ---|---."""
    s = line.strip()
    if "|" not in s and "-" not in s:
        return False
    cells = [c.strip() for c in s.strip("|").split("|")]
    if not cells:
        return False
    return all(c and set(c) <= {"-", ":"} and "-" in c for c in cells)


def _split_row(line: str) -> list[str]:
    """Split a markdown table row into trimmed cells."""
    return [c.strip() for c in line.strip().strip("|").split("|")]


_ORDERED_RE = re.compile(r"^(\d+)\.\s+(.*)$")


def _parse_blocks(content: str) -> list[dict]:
    """
    Parse Markdown content into typed blocks.

    Block types:
      h1/h2/h3   {"type", "text"}
      bullet     {"type", "text"}
      ordered    {"type", "text", "num"}
      hr         {"type"}
      table      {"type", "header": [str], "rows": [[str]]}
      paragraph  {"type", "text"}
    """
    blocks: list[dict] = []
    lines = content.splitlines()
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i].rstrip()
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # ── Table: a row whose next line is a separator ──────────────────────
        if "|" in stripped and i + 1 < n and _is_table_sep(lines[i + 1]):
            header = _split_row(stripped)
            rows: list[list[str]] = []
            i += 2  # skip header + separator
            while i < n and "|" in lines[i] and lines[i].strip():
                if _is_hr(lines[i]):
                    break
                rows.append(_split_row(lines[i]))
                i += 1
            # Normalise ragged rows to header width
            width = len(header)
            rows = [(r + [""] * width)[:width] for r in rows]
            blocks.append({"type": "table", "header": header, "rows": rows})
            continue

        if _is_hr(line):
            blocks.append({"type": "hr"})
        elif stripped.startswith("### "):
            blocks.append({"type": "h3", "text": stripped[4:]})
        elif stripped.startswith("## "):
            blocks.append({"type": "h2", "text": stripped[3:]})
        elif stripped.startswith("# "):
            blocks.append({"type": "h1", "text": stripped[2:]})
        elif stripped.startswith("- ") or stripped.startswith("* "):
            blocks.append({"type": "bullet", "text": stripped[2:]})
        elif _ORDERED_RE.match(stripped):
            m = _ORDERED_RE.match(stripped)
            blocks.append({"type": "ordered", "text": m.group(2), "num": int(m.group(1))})
        else:
            blocks.append({"type": "paragraph", "text": line})
        i += 1

    return blocks


def _dedupe_title(blocks: list[dict], title: str) -> list[dict]:
    """Drop a leading H1 that merely repeats the document title (the cover already shows it)."""
    norm = title.strip().lower()
    for idx, b in enumerate(blocks):
        if b["type"] in ("hr",):
            continue
        if b["type"] == "h1" and _strip_md(b["text"]).strip().lower() == norm:
            return blocks[:idx] + blocks[idx + 1:]
        break
    return blocks


def _strip_md(text: str) -> str:
    """Strip inline **bold**, *italic*, and `code` markers."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text


def _md_to_reportlab(text: str, mono: str = "Courier") -> str:
    """Convert inline **bold** and *italic* to reportlab XML markup."""
    text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text)
    text = re.sub(r"\*(.+?)\*", r"<i>\1</i>", text)
    text = re.sub(r"`(.+?)`", rf"<font name='{mono}'>\1</font>", text)
    # Escape bare ampersands not part of markup
    text = re.sub(r"&(?!amp;|lt;|gt;|quot;|apos;|#)", "&amp;", text)
    return text


# ── Colour helpers ──────────────────────────────────────────────────────────
# A profile gives one accent hex; the generators derive the whole palette from
# it so every format is branded without the profile declaring more than a colour.

_DEFAULT_ACCENT = "#5b6472"   # neutral slate when no profile theme is in context


def _rgb(hex_str: str) -> tuple[float, float, float]:
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return tuple(int(h[i:i + 2], 16) / 255 for i in (0, 2, 4))  # type: ignore[return-value]


def _hex(rgb: tuple[float, float, float]) -> str:
    return "#" + "".join(f"{max(0, min(255, round(c * 255))):02x}" for c in rgb)


def _mix(a: tuple, b: tuple, t: float) -> tuple:
    """Blend a→b by fraction t (0 = a, 1 = b)."""
    return tuple(a[i] * (1 - t) + b[i] * t for i in range(3))


def _luminance(rgb: tuple) -> float:
    r, g, b = rgb
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _readable(hex_str: str) -> str:
    """Darken a bright accent toward black until it reads as body/heading text."""
    rgb = _rgb(hex_str)
    while _luminance(rgb) > 0.42:
        rgb = _mix(rgb, (0.0, 0.0, 0.0), 0.18)
    return _hex(rgb)


def _palette(accent: str) -> dict:
    """Derive the full document palette from a single accent hex."""
    a = _rgb(accent)
    white = (1.0, 1.0, 1.0)
    return {
        "accent": accent,                       # decorative rules / bars (full strength)
        "heading": _readable(accent),           # heading + title-accent text (contrast-safe)
        "header_bg": _hex(_mix(a, white, 0.88)),  # table header fill (light tint)
        "zebra": _hex(_mix(a, white, 0.96)),    # alternating row tint (barely-there)
        "ink": "#1a1a1a",
        "muted": "#6b7280",
        "rule": "#e5e7eb",
    }


def _safe_name(title: str) -> str:
    """Sanitise a title for use in a filename (max 24 chars)."""
    return re.sub(r"[^\w\-]", "_", title)[:24].strip("_")


def _output_path(title: str, ext: str) -> str:
    Path(settings.temp_outputs_dir).mkdir(parents=True, exist_ok=True)
    return os.path.join(
        settings.temp_outputs_dir,
        f"{uuid.uuid4().hex[:8]}_{_safe_name(title)}.{ext}",
    )


# ── Generators ────────────────────────────────────────────────────────────────

def _generate_pptx(title: str, content: str, accent: str = _DEFAULT_ACCENT) -> str:
    from pptx import Presentation        # type: ignore[import]
    from pptx.dml.color import RGBColor  # type: ignore[import]

    heading_rgb = RGBColor.from_string(_palette(accent)["heading"].lstrip("#").upper())

    def _color_title(shape) -> None:
        """Tint a slide's title text with the agent accent."""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = heading_rgb

    prs = Presentation()

    # ── Title slide ──────────────────────────────────────────────────────────
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    _color_title(slide.shapes.title)
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Generated by SPEDA"

    # ── Group blocks into slides at each H2 boundary ─────────────────────────
    slides: list[dict] = []   # [{"title": str, "lines": list[str]}]
    current: dict | None = None

    for block in _dedupe_title(_parse_blocks(content), title):
        kind = block["type"]
        if kind == "h2":
            if current is not None:
                slides.append(current)
            current = {"title": block["text"], "lines": []}
        elif current is None:
            continue
        elif kind == "hr":
            continue
        elif kind == "table":
            current["lines"].append("• " + "  |  ".join(_strip_md(c) for c in block["header"]))
            for row in block["rows"]:
                current["lines"].append("    " + "  |  ".join(_strip_md(c) for c in row))
        else:
            prefix = "• " if kind in ("bullet", "ordered") else ""
            indent = "    " if kind == "h3" else ""
            current["lines"].append(indent + prefix + _strip_md(block["text"]))

    if current is not None:
        slides.append(current)

    # ── Render content slides ─────────────────────────────────────────────────
    content_layout = prs.slide_layouts[1]
    for slide_data in slides:
        sl = prs.slides.add_slide(content_layout)
        sl.shapes.title.text = slide_data["title"]
        _color_title(sl.shapes.title)
        if len(sl.placeholders) > 1:
            tf = sl.placeholders[1].text_frame
            tf.clear()
            for i, line in enumerate(slide_data["lines"]):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    tf.add_paragraph().text = line

    path = _output_path(title, "pptx")
    prs.save(path)
    return path


def _generate_docx(title: str, content: str, accent: str = _DEFAULT_ACCENT) -> str:
    from docx import Document                 # type: ignore[import]
    from docx.oxml import OxmlElement         # type: ignore[import]
    from docx.oxml.ns import qn               # type: ignore[import]
    from docx.shared import RGBColor          # type: ignore[import]

    pal = _palette(accent)
    heading_hex = pal["heading"].lstrip("#")
    header_bg_hex = pal["header_bg"].lstrip("#")
    heading_rgb = RGBColor.from_string(heading_hex.upper())

    def _shade(cell, hex6: str) -> None:
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:fill"), hex6.upper())
        cell._tc.get_or_add_tcPr().append(shd)

    doc = Document()
    doc.add_heading(title, level=0)

    for block in _dedupe_title(_parse_blocks(content), title):
        kind = block["type"]
        if kind == "hr":
            # Horizontal rule rendered as a thin bottom border on an empty paragraph.
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pbdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), heading_hex)
            pbdr.append(bottom)
            pPr.append(pbdr)
            continue
        if kind == "table":
            t = doc.add_table(rows=1, cols=len(block["header"]))
            t.style = "Table Grid"
            for cell, head in zip(t.rows[0].cells, block["header"]):
                cell.text = _strip_md(head)
                _shade(cell, header_bg_hex)
                for run in cell.paragraphs[0].runs:
                    run.bold = True
            for row in block["rows"]:
                cells = t.add_row().cells
                for cell, val in zip(cells, row):
                    cell.text = _strip_md(val)
            continue

        text = _strip_md(block["text"])
        if kind in ("h1", "h2", "h3"):
            h = doc.add_heading(text, level=int(kind[1]))
            for run in h.runs:
                run.font.color.rgb = heading_rgb
        elif kind == "bullet":
            doc.add_paragraph(text, style="List Bullet")
        elif kind == "ordered":
            doc.add_paragraph(text, style="List Number")
        else:
            doc.add_paragraph(text)

    path = _output_path(title, "docx")
    doc.save(path)
    return path


def _generate_pdf(title: str, content: str, accent: str = _DEFAULT_ACCENT) -> str:
    from reportlab.lib import colors                           # type: ignore[import]
    from reportlab.lib.enums import TA_LEFT                     # type: ignore[import]
    from reportlab.lib.pagesizes import A4                      # type: ignore[import]
    from reportlab.lib.styles import (                          # type: ignore[import]
        ParagraphStyle, getSampleStyleSheet,
    )
    from reportlab.lib.units import cm                          # type: ignore[import]
    from reportlab.platypus import (                            # type: ignore[import]
        HRFlowable, ListFlowable, ListItem, Paragraph,
        SimpleDocTemplate, Spacer, Table, TableStyle,
    )

    # ── Palette — derived from the agent's accent ────────────────────────────
    pal = _palette(accent)
    ink = colors.HexColor(pal["ink"])
    accent = colors.HexColor(pal["accent"])
    heading = colors.HexColor(pal["heading"])
    muted = colors.HexColor(pal["muted"])
    rule = colors.HexColor(pal["rule"])
    header_bg = colors.HexColor(pal["header_bg"])
    zebra = colors.HexColor(pal["zebra"])

    # Unicode fonts (Turkish-safe). `md` threads the mono font into inline `code`.
    fonts = _register_pdf_fonts()
    reg, bold, mono = fonts["normal"], fonts["bold"], fonts["mono"]

    def md(text: str) -> str:
        return _md_to_reportlab(text, mono)

    base = getSampleStyleSheet()
    styles = {
        "Title": ParagraphStyle(
            "DocTitle", parent=base["Title"], fontName=bold, fontSize=22, leading=27,
            textColor=ink, spaceAfter=4,
        ),
        "h1": ParagraphStyle(
            "DocH1", parent=base["Heading1"], fontName=bold, fontSize=15, leading=19,
            textColor=ink, spaceBefore=14, spaceAfter=5,
        ),
        "h2": ParagraphStyle(
            "DocH2", parent=base["Heading2"], fontName=bold, fontSize=12.5, leading=16,
            textColor=heading, spaceBefore=11, spaceAfter=4,
        ),
        "h3": ParagraphStyle(
            "DocH3", parent=base["Heading3"], fontName=bold, fontSize=11, leading=14,
            textColor=ink, spaceBefore=8, spaceAfter=3,
        ),
        "body": ParagraphStyle(
            "DocBody", parent=base["Normal"], fontName=reg, fontSize=10, leading=15,
            textColor=ink, alignment=TA_LEFT, spaceAfter=4,
        ),
        "cell": ParagraphStyle(
            "DocCell", parent=base["Normal"], fontName=reg, fontSize=9.5, leading=13,
            textColor=ink,
        ),
        "cellhead": ParagraphStyle(
            "DocCellHead", parent=base["Normal"], fontName=bold, fontSize=9.5, leading=13,
            textColor=ink,
        ),
    }

    path = _output_path(title, "pdf")
    doc = SimpleDocTemplate(
        path, pagesize=A4,
        leftMargin=2.2 * cm, rightMargin=2.2 * cm,
        topMargin=2.2 * cm, bottomMargin=2.2 * cm,
        title=title, author="SPEDA Mark VI",
    )

    story: list = [Paragraph(title, styles["Title"])]
    story.append(HRFlowable(width="100%", thickness=1.2, color=accent,
                            spaceBefore=2, spaceAfter=12))

    # Buffer consecutive bullet/ordered items so they render as one tight list.
    pending: list = []
    pending_kind: str | None = None

    def flush_list():
        nonlocal pending, pending_kind
        if not pending:
            return
        story.append(ListFlowable(
            pending,
            bulletType="bullet" if pending_kind == "bullet" else "1",
            bulletColor=accent, leftIndent=14, bulletFontSize=8,
            spaceBefore=2, spaceAfter=6,
        ))
        pending = []
        pending_kind = None

    for block in _dedupe_title(_parse_blocks(content), title):
        kind = block["type"]

        if kind in ("bullet", "ordered"):
            if pending_kind and pending_kind != kind:
                flush_list()
            pending_kind = kind
            pending.append(ListItem(
                Paragraph(md(block["text"]), styles["body"]),
                leftIndent=14,
            ))
            continue

        flush_list()

        if kind == "hr":
            story.append(HRFlowable(width="100%", thickness=0.6, color=rule,
                                    spaceBefore=8, spaceAfter=8))
        elif kind in ("h1", "h2", "h3"):
            story.append(Paragraph(md(block["text"]), styles[kind]))
        elif kind == "table":
            header = [Paragraph(md(c), styles["cellhead"]) for c in block["header"]]
            data = [header] + [
                [Paragraph(md(c), styles["cell"]) for c in row]
                for row in block["rows"]
            ]
            tbl = Table(data, repeatRows=1, hAlign="LEFT")
            ts = [
                ("BACKGROUND", (0, 0), (-1, 0), header_bg),
                ("LINEBELOW", (0, 0), (-1, 0), 0.8, muted),
                ("GRID", (0, 0), (-1, -1), 0.4, rule),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 7),
                ("RIGHTPADDING", (0, 0), (-1, -1), 7),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]
            for r in range(1, len(data)):
                if r % 2 == 0:
                    ts.append(("BACKGROUND", (0, r), (-1, r), zebra))
            tbl.setStyle(TableStyle(ts))
            story.append(Spacer(1, 0.1 * cm))
            story.append(tbl)
            story.append(Spacer(1, 0.2 * cm))
        else:
            story.append(Paragraph(md(block["text"]), styles["body"]))

    flush_list()
    doc.build(story)
    return path


# ── Skill ─────────────────────────────────────────────────────────────────────

_GENERATORS = {
    "pptx": _generate_pptx,
    "docx": _generate_docx,
    "pdf": _generate_pdf,
}

_REQUIRED_LIBS = {
    "pptx": "python-pptx",
    "docx": "python-docx",
    "pdf": "reportlab",
}


class DocumentsSkill(Skill):
    name = "generate_document"
    description = (
        "Generates a downloadable PPTX, DOCX, or PDF file. "
        "Use ONLY when the user explicitly says they want a file to save, download, print, or send — "
        "e.g. 'create a PDF report', 'make a PowerPoint', 'export as Word'. "
        "NEVER use for flowcharts, diagrams, charts, graphs, dashboards, visualisations, "
        "or any request to 'draw', 'show', 'visualise', or 'render' something — "
        "those are answered with an html or svg fenced code block, not this tool. "
        "Returns the absolute path to the generated file."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "format": {
                "type": "string",
                "enum": ["pptx", "docx", "pdf"],
                "description": "Output file format.",
            },
            "title": {
                "type": "string",
                "description": "Document title (used as the cover title and filename base).",
            },
            "content": {
                "type": "string",
                "description": (
                    "Document body in Markdown. "
                    "For PPTX: use ## headings as slide separators — each H2 becomes a new slide. "
                    "For DOCX/PDF: use # H1 for sections, ## H2 for subsections, - for bullets."
                ),
            },
        },
        "required": ["format", "title", "content"],
    }

    async def execute(self, args: dict, context: AgentContext) -> str:
        fmt = args.get("format", "").lower()
        title = args.get("title", "Untitled")
        content = args.get("content", "")

        logger.info(
            "documents_execute",
            extra={
                "request_id": context.request_id,
                "format": fmt,
                "title": title,
            },
        )

        generator = _GENERATORS.get(fmt)
        if generator is None:
            return f"Unsupported format '{fmt}'. Supported: pptx, docx, pdf."

        # Brand the document with the running agent's accent (Rule 10 — the
        # orchestrator resolves it from the profile onto the context).
        accent = context.extra.get("doc_accent", _DEFAULT_ACCENT)

        try:
            path = generator(title, content, accent)
            logger.info(
                "documents_generated",
                extra={
                    "request_id": context.request_id,
                    "path": path,
                    "format": fmt,
                },
            )
            # Register for delivery — the orchestrator emits a `file` SSE event so
            # the frontend renders a download card.
            from app.core.files import register_file
            meta = register_file(context, path, title=title)
            return (
                f"Created {fmt.upper()} '{title}' ({meta['size']} bytes). "
                f"Delivered to the user as a downloadable file — do NOT paste the "
                f"path or a link, just tell them it's ready."
            )
        except ImportError:
            lib = _REQUIRED_LIBS.get(fmt, "required library")
            return (
                f"Cannot generate {fmt.upper()}: '{lib}' is not installed. "
                f"Run: uv add {lib}"
            )
        except Exception as e:
            logger.error(
                "documents_error",
                extra={
                    "request_id": context.request_id,
                    "format": fmt,
                    "error": str(e),
                },
            )
            return f"Document generation failed: {e}"
