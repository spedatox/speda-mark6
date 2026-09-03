# SPDX-FileCopyrightText: 2026 Ahmet Erol Bayrak
# SPDX-License-Identifier: AGPL-3.0-or-later

import html
import logging
import os
import re
import uuid
from functools import lru_cache
from pathlib import Path

from app.config import settings
from app.core.context import AgentContext
from app.skills.base import Skill

logger = logging.getLogger(__name__)

# ── PDF fonts ───────────────────────────────────────────────────────────────
# The DejaVu Sans family (full Unicode, permissive licence) covers the
# Turkish-specific letters ğ ş ı İ that many system font fallbacks lack. Each
# generated PDF's CSS declares @font-face against these bundled files (via
# file:// URIs — see _pdf_font_faces) rather than relying on whatever fonts the
# container image happens to ship, so rendering is identical on every
# deployment regardless of host font packages.
_FONTS_DIR = Path(__file__).parent / "fonts"

_PDF_FONT_VARIANTS = [
    ("DejaVu Sans", "normal", "normal", "DejaVuSans.ttf"),
    ("DejaVu Sans", "bold", "normal", "DejaVuSans-Bold.ttf"),
    ("DejaVu Sans", "normal", "italic", "DejaVuSans-Oblique.ttf"),
    ("DejaVu Sans", "bold", "italic", "DejaVuSans-BoldOblique.ttf"),
    ("DejaVu Sans Mono", "normal", "normal", "DejaVuSansMono.ttf"),
]


@lru_cache(maxsize=1)
def _pdf_font_faces() -> str:
    """Build the @font-face CSS block for the bundled DejaVu family. A missing
    file is skipped (logged) rather than raising — the document still renders,
    just falls back to whatever font WeasyPrint substitutes."""
    faces: list[str] = []
    for family, weight, style, fname in _PDF_FONT_VARIANTS:
        path = _FONTS_DIR / fname
        if not path.exists():
            logger.warning("pdf_font_missing", extra={"file": fname})
            continue
        faces.append(
            f"@font-face {{ font-family: '{family}'; font-weight: {weight}; "
            f"font-style: {style}; src: url('{path.resolve().as_uri()}'); }}"
        )
    return "\n".join(faces)


# ── Markdown parsing ──────────────────────────────────────────────────────────

def _is_hr(line: str) -> bool:
    """A thematic break: 3+ of -, *, or _ (no pipes, so table separators are excluded)."""
    s = line.strip().replace(" ", "")
    return len(s) >= 3 and (set(s) == {"-"} or set(s) == {"*"} or set(s) == {"_"})


def _is_table_sep(line: str) -> bool:
    """A table separator row, e.g. |---|:--:|---| or ---|---."""
    s = line.strip()
    if "|" not in s and "-" not in s:
        return False
    cells = [c.strip() for c in s.strip("|").split("|")]
    if not cells:
        return False
    return all(c and set(c) <= {"-", ":"} and "-" in c for c in cells)


def _split_row(line: str) -> list[str]:
    """Split a markdown table row into trimmed cells."""
    return [c.strip() for c in line.strip().strip("|").split("|")]


_ORDERED_RE = re.compile(r"^(\d+)\.\s+(.*)$")


def _parse_blocks(content: str) -> list[dict]:
    """
    Parse Markdown content into typed blocks.

    Block types:
      h1/h2/h3   {"type", "text"}
      bullet     {"type", "text"}
      ordered    {"type", "text", "num"}
      hr         {"type"}
      table      {"type", "header": [str], "rows": [[str]]}
      paragraph  {"type", "text"}
    """
    blocks: list[dict] = []
    lines = content.splitlines()
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i].rstrip()
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # ── Table: a row whose next line is a separator ──────────────────────
        if "|" in stripped and i + 1 < n and _is_table_sep(lines[i + 1]):
            header = _split_row(stripped)
            rows: list[list[str]] = []
            i += 2  # skip header + separator
            while i < n and "|" in lines[i] and lines[i].strip():
                if _is_hr(lines[i]):
                    break
                rows.append(_split_row(lines[i]))
                i += 1
            # Normalise ragged rows to header width
            width = len(header)
            rows = [(r + [""] * width)[:width] for r in rows]
            blocks.append({"type": "table", "header": header, "rows": rows})
            continue

        if _is_hr(line):
            blocks.append({"type": "hr"})
        elif stripped.startswith("### "):
            blocks.append({"type": "h3", "text": stripped[4:]})
        elif stripped.startswith("## "):
            blocks.append({"type": "h2", "text": stripped[3:]})
        elif stripped.startswith("# "):
            blocks.append({"type": "h1", "text": stripped[2:]})
        elif stripped.startswith("- ") or stripped.startswith("* "):
            blocks.append({"type": "bullet", "text": stripped[2:]})
        elif _ORDERED_RE.match(stripped):
            m = _ORDERED_RE.match(stripped)
            blocks.append({"type": "ordered", "text": m.group(2), "num": int(m.group(1))})
        else:
            blocks.append({"type": "paragraph", "text": line})
        i += 1

    return blocks


def _dedupe_title(blocks: list[dict], title: str) -> list[dict]:
    """Drop a leading H1 that merely repeats the document title (the cover already shows it)."""
    norm = title.strip().lower()
    for idx, b in enumerate(blocks):
        if b["type"] in ("hr",):
            continue
        if b["type"] == "h1" and _strip_md(b["text"]).strip().lower() == norm:
            return blocks[:idx] + blocks[idx + 1:]
        break
    return blocks


def _strip_md(text: str) -> str:
    """Strip inline **bold**, *italic*, and `code` markers."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text


def _md_to_html(text: str) -> str:
    """Escape text for safe HTML embedding, then convert inline **bold**,
    *italic*, and `code` markers to their tag equivalents."""
    text = html.escape(text, quote=False)
    text = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
    text = re.sub(r"\*(.+?)\*", r"<em>\1</em>", text)
    text = re.sub(r"`(.+?)`", r"<code>\1</code>", text)
    return text


# ── Colour helpers ──────────────────────────────────────────────────────────
# A profile gives one accent hex; the generators derive the whole palette from
# it so every format is branded without the profile declaring more than a colour.

_DEFAULT_ACCENT = "#5b6472"   # neutral slate when no profile theme is in context

# ── Signature ────────────────────────────────────────────────────────────────
# Every artifact the system leaves behind carries the name and mark of the agent
# that made it (prompts/core/13_signature.md) — "Generated by Sentinel Mark II".
# The author string arrives already composed as profile.signed_name; for
# generated files the stamp is applied here rather than trusted to the model
# (the agents are told NOT to write a "Generated by" line into the Markdown),
# so the footer appears exactly once, in the right place, in the right type.

_DEFAULT_AUTHOR = "Speda Mark VI"   # unattributed runs still carry a mark


def _signature(author: str | None) -> str:
    """`Generated by Sentinel Mark II` — the agent's own name and iteration."""
    return f"Generated by {(author or '').strip() or _DEFAULT_AUTHOR}"


def _rgb(hex_str: str) -> tuple[float, float, float]:
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return tuple(int(h[i:i + 2], 16) / 255 for i in (0, 2, 4))  # type: ignore[return-value]


def _hex(rgb: tuple[float, float, float]) -> str:
    return "#" + "".join(f"{max(0, min(255, round(c * 255))):02x}" for c in rgb)


def _mix(a: tuple, b: tuple, t: float) -> tuple:
    """Blend a→b by fraction t (0 = a, 1 = b)."""
    return tuple(a[i] * (1 - t) + b[i] * t for i in range(3))


def _luminance(rgb: tuple) -> float:
    r, g, b = rgb
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _readable(hex_str: str) -> str:
    """Darken a bright accent toward black until it reads as body/heading text."""
    rgb = _rgb(hex_str)
    while _luminance(rgb) > 0.42:
        rgb = _mix(rgb, (0.0, 0.0, 0.0), 0.18)
    return _hex(rgb)


def _palette(accent: str) -> dict:
    """Derive the full document palette from a single accent hex."""
    a = _rgb(accent)
    white = (1.0, 1.0, 1.0)
    return {
        "accent": accent,                       # decorative rules / bars (full strength)
        "heading": _readable(accent),           # heading + title-accent text (contrast-safe)
        "header_bg": _hex(_mix(a, white, 0.88)),  # table header fill (light tint)
        "zebra": _hex(_mix(a, white, 0.96)),    # alternating row tint (barely-there)
        "ink": "#1a1a1a",
        "muted": "#6b7280",
        "rule": "#e5e7eb",
    }


def _safe_name(title: str) -> str:
    """Sanitise a title for use in a filename (max 24 chars)."""
    return re.sub(r"[^\w\-]", "_", title)[:24].strip("_")


def _output_path(title: str, ext: str) -> str:
    Path(settings.temp_outputs_dir).mkdir(parents=True, exist_ok=True)
    return os.path.join(
        settings.temp_outputs_dir,
        f"{uuid.uuid4().hex[:8]}_{_safe_name(title)}.{ext}",
    )


# ── Generators ────────────────────────────────────────────────────────────────

def _generate_pptx(title: str, content: str, accent: str = _DEFAULT_ACCENT,
                   author: str | None = None) -> str:
    from pptx import Presentation        # type: ignore[import]
    from pptx.dml.color import RGBColor  # type: ignore[import]

    heading_rgb = RGBColor.from_string(_palette(accent)["heading"].lstrip("#").upper())

    def _color_title(shape) -> None:
        """Tint a slide's title text with the agent accent."""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = heading_rgb

    prs = Presentation()

    # ── Title slide ──────────────────────────────────────────────────────────
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    _color_title(slide.shapes.title)
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = _signature(author)

    # ── Group blocks into slides at each H2 boundary ─────────────────────────
    slides: list[dict] = []   # [{"title": str, "lines": list[str]}]
    current: dict | None = None

    for block in _dedupe_title(_parse_blocks(content), title):
        kind = block["type"]
        if kind == "h2":
            if current is not None:
                slides.append(current)
            current = {"title": block["text"], "lines": []}
        elif current is None:
            continue
        elif kind == "hr":
            continue
        elif kind == "table":
            current["lines"].append("• " + "  |  ".join(_strip_md(c) for c in block["header"]))
            for row in block["rows"]:
                current["lines"].append("    " + "  |  ".join(_strip_md(c) for c in row))
        else:
            prefix = "• " if kind in ("bullet", "ordered") else ""
            indent = "    " if kind == "h3" else ""
            current["lines"].append(indent + prefix + _strip_md(block["text"]))

    if current is not None:
        slides.append(current)

    # ── Render content slides ─────────────────────────────────────────────────
    content_layout = prs.slide_layouts[1]
    for slide_data in slides:
        sl = prs.slides.add_slide(content_layout)
        sl.shapes.title.text = slide_data["title"]
        _color_title(sl.shapes.title)
        if len(sl.placeholders) > 1:
            tf = sl.placeholders[1].text_frame
            tf.clear()
            for i, line in enumerate(slide_data["lines"]):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    tf.add_paragraph().text = line

    path = _output_path(title, "pptx")
    prs.save(path)
    return path


def _generate_docx(title: str, content: str, accent: str = _DEFAULT_ACCENT,
                   author: str | None = None) -> str:
    from docx import Document                 # type: ignore[import]
    from docx.oxml import OxmlElement         # type: ignore[import]
    from docx.oxml.ns import qn               # type: ignore[import]
    from docx.shared import Pt, RGBColor      # type: ignore[import]

    pal = _palette(accent)
    heading_hex = pal["heading"].lstrip("#")
    header_bg_hex = pal["header_bg"].lstrip("#")
    heading_rgb = RGBColor.from_string(heading_hex.upper())

    def _shade(cell, hex6: str) -> None:
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:fill"), hex6.upper())
        cell._tc.get_or_add_tcPr().append(shd)

    doc = Document()
    doc.add_heading(title, level=0)

    for block in _dedupe_title(_parse_blocks(content), title):
        kind = block["type"]
        if kind == "hr":
            # Horizontal rule rendered as a thin bottom border on an empty paragraph.
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pbdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), heading_hex)
            pbdr.append(bottom)
            pPr.append(pbdr)
            continue
        if kind == "table":
            t = doc.add_table(rows=1, cols=len(block["header"]))
            t.style = "Table Grid"
            for cell, head in zip(t.rows[0].cells, block["header"]):
                cell.text = _strip_md(head)
                _shade(cell, header_bg_hex)
                for run in cell.paragraphs[0].runs:
                    run.bold = True
            for row in block["rows"]:
                cells = t.add_row().cells
                for cell, val in zip(cells, row):
                    cell.text = _strip_md(val)
            continue

        text = _strip_md(block["text"])
        if kind in ("h1", "h2", "h3"):
            h = doc.add_heading(text, level=int(kind[1]))
            for run in h.runs:
                run.font.color.rgb = heading_rgb
        elif kind == "bullet":
            doc.add_paragraph(text, style="List Bullet")
        elif kind == "ordered":
            doc.add_paragraph(text, style="List Number")
        else:
            doc.add_paragraph(text)

    # Signature — last thing in the document, muted and smaller than the body.
    sig = doc.add_paragraph()
    sig_run = sig.add_run(_signature(author))
    sig_run.italic = True
    sig_run.font.size = Pt(8)
    sig_run.font.color.rgb = RGBColor.from_string(pal["muted"].lstrip("#").upper())

    path = _output_path(title, "docx")
    doc.save(path)
    return path


def _generate_pdf(title: str, content: str, accent: str = _DEFAULT_ACCENT,
                  author: str | None = None) -> str:
    from weasyprint import HTML  # type: ignore[import]

    pal = _palette(accent)

    # Buffer consecutive bullet/ordered items so they render as one tight list.
    body: list[str] = []
    pending: list[str] = []
    pending_kind: str | None = None

    def flush_list():
        nonlocal pending, pending_kind
        if not pending:
            return
        tag = "ul" if pending_kind == "bullet" else "ol"
        body.append(f"<{tag}>" + "".join(pending) + f"</{tag}>")
        pending = []
        pending_kind = None

    for block in _dedupe_title(_parse_blocks(content), title):
        kind = block["type"]

        if kind in ("bullet", "ordered"):
            if pending_kind and pending_kind != kind:
                flush_list()
            pending_kind = kind
            pending.append(f"<li>{_md_to_html(block['text'])}</li>")
            continue

        flush_list()

        if kind == "hr":
            body.append('<hr class="thin">')
        elif kind in ("h1", "h2", "h3"):
            body.append(f"<{kind}>{_md_to_html(block['text'])}</{kind}>")
        elif kind == "table":
            head_cells = "".join(f"<th>{_md_to_html(c)}</th>" for c in block["header"])
            row_html = "".join(
                "<tr>" + "".join(f"<td>{_md_to_html(c)}</td>" for c in row) + "</tr>"
                for row in block["rows"]
            )
            body.append(f"<table><thead><tr>{head_cells}</tr></thead><tbody>{row_html}</tbody></table>")
        else:
            body.append(f"<p>{_md_to_html(block['text'])}</p>")

    flush_list()

    doc_html = f"""<!DOCTYPE html>
<html lang="tr"><head><meta charset="utf-8"><title>{html.escape(title)}</title><style>
{_pdf_font_faces()}
@page {{ size: A4; margin: 2.2cm; }}
* {{ box-sizing: border-box; }}
body {{ font-family: 'DejaVu Sans', sans-serif; color: {pal['ink']}; font-size: 10.5pt; line-height: 1.5; margin: 0; }}
h1.doc-title {{ font-weight: bold; font-size: 22pt; margin: 0 0 4pt; color: {pal['ink']}; }}
hr.title-rule {{ border: none; border-top: 1.2pt solid {pal['accent']}; margin: 2pt 0 14pt; }}
h1 {{ font-size: 15pt; margin: 14pt 0 5pt; color: {pal['ink']}; }}
h2 {{ font-size: 12.5pt; margin: 11pt 0 4pt; color: {pal['heading']}; }}
h3 {{ font-size: 11pt; margin: 8pt 0 3pt; color: {pal['ink']}; }}
p {{ margin: 0 0 4pt; }}
ul, ol {{ margin: 2pt 0 6pt 0; padding-left: 14pt; }}
li {{ margin-bottom: 2pt; }}
li::marker {{ color: {pal['accent']}; }}
hr.thin {{ border: none; border-top: 0.6pt solid {pal['rule']}; margin: 8pt 0; }}
table {{ width: 100%; border-collapse: collapse; margin: 6pt 0 8pt; font-size: 9.5pt; }}
th, td {{ border: 0.4pt solid {pal['rule']}; padding: 5pt 7pt; text-align: left; vertical-align: middle; }}
th {{ background: {pal['header_bg']}; border-bottom: 0.8pt solid {pal['muted']}; font-weight: bold; }}
tbody tr:nth-child(even) td {{ background: {pal['zebra']}; }}
code {{ font-family: 'DejaVu Sans Mono', monospace; }}
p.signature {{ margin: 18pt 0 0; padding-top: 6pt; border-top: 0.4pt solid {pal['rule']};
               color: {pal['muted']}; font-size: 8pt; font-style: italic; }}
</style></head>
<body>
<h1 class="doc-title">{html.escape(title)}</h1>
<hr class="title-rule">
{''.join(body)}
<p class="signature">{html.escape(_signature(author))}</p>
</body></html>"""

    path = _output_path(title, "pdf")
    HTML(string=doc_html).write_pdf(path)
    return path


# ── Skill ─────────────────────────────────────────────────────────────────────

_GENERATORS = {
    "pptx": _generate_pptx,
    "docx": _generate_docx,
    "pdf": _generate_pdf,
}

_REQUIRED_LIBS = {
    "pptx": "python-pptx",
    "docx": "python-docx",
    "pdf": "weasyprint",
}


class DocumentsSkill(Skill):
    name = "generate_document"
    deferred = True
    search_keywords = "powerpoint pptx slides deck word docx pdf report document presentation export write file"
    description = (
        "Generates a downloadable PPTX, DOCX, or PDF file. "
        "Use ONLY when the user explicitly says they want a file to save, download, print, or send — "
        "e.g. 'create a PDF report', 'make a PowerPoint', 'export as Word'. "
        "NEVER use for flowcharts, diagrams, charts, graphs, dashboards, visualisations, "
        "or any request to 'draw', 'show', 'visualise', or 'render' something — "
        "those are answered with an html or svg fenced code block, not this tool. "
        "Returns the absolute path to the generated file."
    )
    input_schema = {
        "type": "object",
        "properties": {
            "format": {
                "type": "string",
                "enum": ["pptx", "docx", "pdf"],
                "description": "Output file format.",
            },
            "title": {
                "type": "string",
                "description": "Document title (used as the cover title and filename base).",
            },
            "content": {
                "type": "string",
                "description": (
                    "Document body in Markdown. "
                    "For PPTX: use ## headings as slide separators — each H2 becomes a new slide. "
                    "For DOCX/PDF: use # H1 for sections, ## H2 for subsections, - for bullets."
                ),
            },
        },
        "required": ["format", "title", "content"],
    }

    async def execute(self, args: dict, context: AgentContext) -> str:
        fmt = args.get("format", "").lower()
        title = args.get("title", "Untitled")
        content = args.get("content", "")

        logger.info(
            "documents_execute",
            extra={
                "request_id": context.request_id,
                "format": fmt,
                "title": title,
            },
        )

        generator = _GENERATORS.get(fmt)
        if generator is None:
            return f"Unsupported format '{fmt}'. Supported: pptx, docx, pdf."

        # Brand the document with the running agent's accent (Rule 10 — the
        # orchestrator resolves it from the profile onto the context).
        accent = context.extra.get("doc_accent", _DEFAULT_ACCENT)
        # …and sign it with that agent's name (prompts/core/13_signature.md).
        author = context.extra.get("doc_author", _DEFAULT_AUTHOR)

        try:
            path = generator(title, content, accent, author)
            logger.info(
                "documents_generated",
                extra={
                    "request_id": context.request_id,
                    "path": path,
                    "format": fmt,
                },
            )
            # Register for delivery — the orchestrator emits a `file` SSE event so
            # the frontend renders a download card.
            from app.core.files import register_file
            meta = register_file(context, path, title=title)
            return (
                f"Created {fmt.upper()} '{title}' ({meta['size']} bytes). "
                f"The file is named '{meta['name']}'. "
                f"Delivered to the user as a downloadable file — do NOT paste the "
                f"path or a link, just tell them it's ready. "
                f"If sending via Telegram, use filename '{meta['name']}'."
            )
        except ImportError:
            lib = _REQUIRED_LIBS.get(fmt, "required library")
            return (
                f"Cannot generate {fmt.upper()}: '{lib}' is not installed. "
                f"Run: uv add {lib}"
            )
        except Exception as e:
            logger.error(
                "documents_error",
                extra={
                    "request_id": context.request_id,
                    "format": fmt,
                    "error": str(e),
                },
            )
            return f"Document generation failed: {e}"
