"""
Attachment text extraction — turns an uploaded non-image file into plain text
that can be embedded in the user turn.

Why text, not native document blocks: only Anthropic accepts a `document`
content block, and even there only PDF/plain-text. The backend serves six
providers (Anthropic, OpenAI, Gemini, z.ai, DeepSeek, Ollama), several of which
have zero document support (open-weight GLM/DeepSeek, local Ollama). Extracting
text here and injecting it as an ordinary text block is the one representation
that reaches the model IDENTICALLY on every provider — it survives the
Anthropic→chat-completions translation in llm_client untouched, because it is
just text. It also keeps the document in context across turns (the block is
persisted and re-sent), and works in the Dead Zone where only local weights run.

Images are NOT handled here — they keep their native vision path (image blocks)
in the chat router, which every provider's translation layer already supports.

Extraction is best-effort and never raises: a corrupt or unsupported file yields
a short note in place of its text, so one bad attachment never fails the turn.
"""

import base64
import csv
import io
import logging

logger = logging.getLogger(__name__)

# Per-file cap on extracted characters. A large document would otherwise blow
# the context window and per-turn cost (~4 chars/token → 120k chars ≈ 30k
# tokens). Truncated files are marked so the model knows it saw only a prefix.
_MAX_CHARS = 120_000

# Media types (and filename suffixes) we treat as already-plaintext: decode the
# bytes as UTF-8 directly, no library needed. Covers code, config, data and
# markup formats a user is likely to drop into chat.
_TEXT_SUFFIXES = (
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl", ".yaml",
    ".yml", ".xml", ".html", ".htm", ".css", ".ini", ".toml", ".cfg", ".conf",
    ".log", ".rtf", ".tex", ".srt", ".vtt",
    # source code
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".h", ".cpp", ".cc",
    ".hpp", ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala",
    ".sh", ".bash", ".zsh", ".ps1", ".sql", ".r", ".m", ".pl", ".lua", ".dart",
    ".vue", ".svelte", ".gradle", ".dockerfile", ".env", ".gitignore",
)


def _is_text_like(name: str, media_type: str) -> bool:
    mt = (media_type or "").lower()
    if mt.startswith("text/"):
        return True
    if mt in ("application/json", "application/xml", "application/x-yaml", "application/javascript"):
        return True
    return name.lower().endswith(_TEXT_SUFFIXES)


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader  # type: ignore[import]

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(f"--- Page {i} ---\n{text}")
    return "\n\n".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document  # type: ignore[import]

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # Tables carry a lot of the meaning in real-world .docx — flatten each row.
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        lines.append(txt)
        if lines:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(lines))
    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook  # type: ignore[import]

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                rows.append(",".join(cells))
        if rows:
            parts.append(f"--- Sheet: {ws.title} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _extract_csv(data: bytes) -> str:
    text = _decode_text(data)
    # Normalise into a clean, aligned-ish table the model reads well; falls back
    # to the raw text if the dialect can't be sniffed.
    try:
        reader = csv.reader(io.StringIO(text))
        rows = [" | ".join(cell.strip() for cell in row) for row in reader]
        return "\n".join(rows)
    except csv.Error:
        return text


# media_type / suffix → extractor. Checked in order of specificity.
_PDF_TYPES = ("application/pdf",)
_DOCX_TYPES = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
)
_PPTX_TYPES = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
)
_XLSX_TYPES = (
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
)


def _raw_extract(name: str, media_type: str, data: bytes) -> str:
    n = name.lower()
    mt = (media_type or "").lower()

    if mt in _PDF_TYPES or n.endswith(".pdf"):
        return _extract_pdf(data)
    if mt in _DOCX_TYPES or n.endswith(".docx"):
        return _extract_docx(data)
    if mt in _PPTX_TYPES or n.endswith(".pptx"):
        return _extract_pptx(data)
    if mt in _XLSX_TYPES or n.endswith((".xlsx", ".xlsm")):
        return _extract_xlsx(data)
    if n.endswith((".csv", ".tsv")) or mt == "text/csv":
        return _extract_csv(data)
    if _is_text_like(name, media_type):
        return _decode_text(data)
    # Unknown/binary — last-ditch UTF-8 decode; if it's real binary this yields
    # mostly replacement chars, so guard against emitting garbage.
    text = _decode_text(data)
    printable = sum(c.isprintable() or c.isspace() for c in text[:2000])
    if text and printable / max(1, len(text[:2000])) > 0.85:
        return text
    return ""


def extract_text(name: str, media_type: str, data_b64: str) -> str:
    """
    Extract the text of one attachment and return it wrapped in a labelled
    envelope ready to drop into the user turn as a text block. Best-effort:
    never raises — on any failure it returns a short note naming the file so the
    model at least knows an attachment was present.
    """
    try:
        data = base64.b64decode(data_b64)
    except Exception:
        return f"[Attachment '{name}' could not be decoded and was skipped.]"

    try:
        body = _raw_extract(name, media_type, data)
    except ImportError as exc:
        logger.warning("attachment_extract_missing_lib", extra={"name": name, "error": str(exc)})
        return f"[Attachment '{name}' ({media_type}): text extraction unavailable on the server.]"
    except Exception as exc:
        logger.warning("attachment_extract_failed", extra={"name": name, "error": str(exc)})
        return f"[Attachment '{name}' ({media_type}) could not be read.]"

    body = (body or "").strip()
    if not body:
        return (
            f"[Attachment '{name}' ({media_type}) contained no extractable text "
            f"(it may be a scanned image or an unsupported binary format).]"
        )

    truncated = ""
    if len(body) > _MAX_CHARS:
        body = body[:_MAX_CHARS]
        truncated = f"\n\n[... '{name}' truncated at {_MAX_CHARS:,} characters ...]"

    return (
        f"[Attached file: {name} ({media_type})]\n"
        f"{body}{truncated}\n"
        f"[End of {name}]"
    )


def build_user_content(
    message: str,
    attachments: list,
    documents: list,
) -> list | str:
    """Assemble one user turn's content blocks from a ChatRequest's message +
    attachments/documents. Images keep their native vision blocks (every
    provider's translation layer supports them). Non-image files are extracted
    to plain text here and embedded as text blocks, so a PDF/DOCX/XLSX/CSV
    reaches the model IDENTICALLY on every provider — including open-weight and
    local ones with no native document support. Returns the bare message string
    when there are no attachments/documents (no meta block needed)."""
    if not attachments and not documents:
        return message

    blocks: list = [
        {
            "type": "image",
            "source": {"type": "base64", "media_type": att.media_type, "data": att.data},
        }
        for att in attachments
    ]
    for doc in documents:
        blocks.append({"type": "text", "text": extract_text(doc.name, doc.media_type, doc.data)})
    if message:
        blocks.append({"type": "text", "text": message})
    # Display-only meta: upload chips on the user bubble survive a reload.
    # Stripped from the history sent to the model by SessionManager._clean.
    if documents:
        # `text` carries the user's own message so the reloaded bubble shows it
        # — NOT the wall of extracted document text (which lives in real text
        # blocks the model reads but the UI must not echo).
        blocks.append({
            "type": "_speda_meta",
            "uploads": [{"name": d.name, "size": d.size or 0} for d in documents],
            "text": message or "",
        })
    return blocks
